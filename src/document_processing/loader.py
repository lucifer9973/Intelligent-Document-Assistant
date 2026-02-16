"""
Document Loader - Handles reading various document formats
Supports: PDF, DOCX, TXT, PPTX
"""
import os
from typing import List, Optional
from dataclasses import dataclass
import logging

logger = logging.getLogger(__name__)


@dataclass
class Document:
    """Represents a loaded document"""
    content: str
    filename: str
    format: str
    metadata: dict


class DocumentLoader:
    """Load and extract text from various document formats"""
    
    def __init__(self, supported_formats: List[str] = None):
        """
        Initialize DocumentLoader
        
        Args:
            supported_formats: List of supported file formats (pdf, docx, txt, pptx)
        """
        self.supported_formats = supported_formats or ["pdf", "docx", "txt", "pptx"]
    
    def load(self, file_path: str) -> Optional[Document]:
        """
        Load document from file path
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Document object or None if loading fails
        """
        if not os.path.exists(file_path):
            logger.error(f"File not found: {file_path}")
            return None
        
        file_ext = os.path.splitext(file_path)[1].lower().strip(".")
        filename = os.path.basename(file_path)
        
        if file_ext not in self.supported_formats:
            logger.error(f"Unsupported format: {file_ext}")
            return None
        
        try:
            if file_ext == "pdf":
                content = self._load_pdf(file_path)
            elif file_ext == "docx":
                content = self._load_docx(file_path)
            elif file_ext == "txt":
                content = self._load_txt(file_path)
            elif file_ext == "pptx":
                content = self._load_pptx(file_path)
            else:
                return None
            
            return Document(
                content=content,
                filename=filename,
                format=file_ext,
                metadata={
                    "file_path": file_path,
                    "file_size": os.path.getsize(file_path),
                }
            )
        except Exception as e:
            logger.error(f"Error loading document {filename}: {str(e)}")
            return None
    
    def _load_pdf(self, file_path: str) -> str:
        """Extract text from PDF"""
        try:
            import PyPDF2
            content = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    content += page.extract_text()
            return content
        except ImportError:
            logger.error("PyPDF2 not installed")
            return ""
    
    def _load_docx(self, file_path: str) -> str:
        """Extract text from DOCX"""
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(file_path)
            content = "\n".join([para.text for para in doc.paragraphs])
            return content
        except ImportError:
            logger.error("python-docx not installed")
            return ""
    
    def _load_txt(self, file_path: str) -> str:
        """Extract text from TXT"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            logger.error(f"Error reading text file: {str(e)}")
            return ""
    
    def _load_pptx(self, file_path: str) -> str:
        """Extract text from PPTX"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            content = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
            return content
        except ImportError:
            logger.error("python-pptx not installed")
            return ""
